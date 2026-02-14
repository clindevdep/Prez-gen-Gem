import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Zentiva Brand Colors
ZENTIVA_BLUE = RGBColor(12, 27, 125)   # #0C1B7D
ZENTIVA_TEAL = RGBColor(0, 169, 143)   # #00A98F

def get_layout_by_name(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[1]

def remove_all_slides(prs):
    """Removes all slides from the presentation object in a way that avoids XML corruption."""
    rIds = [prs.slides._sldIdLst[i].rId for i in range(len(prs.slides))]
    for rId in rIds:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.clear()

def create_presentation(title, slides_content, output_path, template_path=None):
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
        remove_all_slides(prs)
    else:
        prs = Presentation()

    # Title Slide
    title_layout = get_layout_by_name(prs, "Title Slide")
    slide = prs.slides.add_slide(title_layout)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title
        if shape.placeholder_format.idx == 1:
            shape.text = "Strategic Market Insights v003"

    # Content Slides
    for slide_data in slides_content:
        layout_name = slide_data.get('layout', "2_Title and Content")
        layout = get_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']
        
        for shape in slide.placeholders:
            ph = shape.placeholder_format
            if ph.idx == 14 and 'content' in slide_data:
                shape.text = slide_data['content']
            elif ph.idx == 15 and 'content2' in slide_data:
                shape.text = slide_data['content2']
            elif ph.type == 18 and 'image' in slide_data: # PICTURE
                if os.path.exists(slide_data['image']):
                    shape.insert_picture(slide_data['image'])
                else:
                    print(f"Warning: Image {slide_data['image']} not found.")

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python generate_pptx.py <title> <output_path>")
        sys.exit(1)

    title = sys.argv[1]
    output_path = sys.argv[2]
    
    script_dir = os.path.dirname(os.path.abspath(__file__))
    assets_dir = os.path.join(script_dir, "..", "assets")
    
    # Structure for v003
    slides = [
        {
            "title": "Market Overview & Growth",
            "content": "The global biosimilars market is projected to reach ~$49B by 2026.\n\nKey growth engine for Zentiva's portfolio.",
            "layout": "2_Title and Content"
        },
        {
            "title": "Dual Impact Analysis",
            "content": "Biosimilars drive cost savings (> $30B/year) and improve system sustainability.",
            "content2": "Patent cliffs create 'biosimilar voids' for generics to strategically fill.",
            "layout": "17_Title and Content" 
        },
        {
            "title": "Growth Trajectory",
            "content": "Visual representation of market expansion from 2021-2026.",
            "image": os.path.join(assets_dir, "market_chart.png"), 
            "layout": "10_Title and Content" 
        }
    ]
    
    template = os.path.join(assets_dir, "Brand_v001.pptx")
    
    create_presentation(title, slides, output_path, template if os.path.exists(template) else None)
