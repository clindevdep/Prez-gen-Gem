from pptx import Presentation

def audit_template(path):
    prs = Presentation(path)
    print(f"--- Template Audit: {path} ---")
    
    print("\n[Slide Layouts]")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Index {i}: {layout.name}")
        for shape in layout.shapes:
            if shape.is_placeholder:
                print(f"  - Placeholder: {shape.placeholder_format.type} (idx {shape.placeholder_format.idx})")

    print("\n[Color Samples from Master Shapes]")
    for shape in prs.slide_master.shapes:
        try:
            if hasattr(shape, "fill") and hasattr(shape.fill, "fore_color"):
                # 1 = RGB, 2 = Theme
                color_type = shape.fill.fore_color.type
                if color_type == 1:
                    print(f"Shape: {shape.name} | RGB: {shape.fill.fore_color.rgb}")
                elif color_type == 2:
                    print(f"Shape: {shape.name} | Theme Color Index: {shape.fill.fore_color.theme_color}")
        except:
            continue

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        audit_template(sys.argv[1])
